
import os
import sys
from pptx import Presentation
from pptx.slide import Slides

# print(len(sys.argv))
# print(sys.argv[1])
# print(os.path.isfile(sys.argv[1]))

if len(sys.argv) != 2 or not(os.path.isfile(sys.argv[1])): 
    sys.exit("Ingresa el nombre del archivo pptx")
        
summary_num_slides = [114, 99, 50, 67, 81, 88, 94, 33, 23]

# https://stackoverflow.com/questions/37375687/python-to-remove-copy-ppt-slides
def dropSlides(slidesToKeep, prs):
    indexesToRemove = [x for x in range(1, len(prs.slides._sldIdLst)+1) if x not in slidesToKeep]

    for i, slide in enumerate(prs.slides):
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}

        if i+1 in indexesToRemove:
            slide_id = slide.slide_id

            prs.part.drop_rel(id_dict[slide_id][1])
            del prs.slides._sldIdLst[id_dict[slide_id][0]]

    return prs


pptx_file = sys.argv[1] # 'QPM Corrimiento 2024-11 Coeficiente autorregresivo.pptx'
summary_pptx_file = '[Summary] - ' + os.path.basename(pptx_file)

print('Reading original file...')
ppt = Presentation(pptx_file)
summary_ppt = dropSlides(summary_num_slides, ppt)

print('Writing summary file...')
summary_ppt.save(summary_pptx_file)
print('[done]')